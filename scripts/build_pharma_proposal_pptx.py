"""Build a clean, professional Fratelanza Pharma Field Force proposal from scratch."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent.parent / "Fratelanza_Pharma_Field_Force_Technical_Proposal_FIXED.pptx"

C_PRIMARY = RGBColor(0x16, 0x69, 0xB4)
C_ACCENT = RGBColor(0x1E, 0xAE, 0xD7)
C_TEAL = RGBColor(0x19, 0x91, 0x82)
C_GOLD = RGBColor(0xEB, 0x8C, 0x2D)
C_DARK = RGBColor(0x0F, 0x2A, 0x44)
C_TEXT = RGBColor(0x1E, 0x2D, 0x3D)
C_MUTED = RGBColor(0x5C, 0x6B, 0x7A)
C_CARD = RGBColor(0xF7, 0xFA, 0xFC)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_ROW = RGBColor(0xED, 0xF2, 0xF7)
C_BORDER = RGBColor(0xD5, 0xDF, 0xEA)
FONT = "Segoe UI"

ACCENTS = [C_PRIMARY, C_ACCENT, C_TEAL, C_GOLD]


def font(run, pt, *, bold=False, color=C_TEXT):
    run.font.name = FONT
    run.font.size = Pt(pt)
    run.font.bold = bold
    run.font.color.rgb = color


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, l, t, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(0.5)
    else:
        s.line.fill.background()
    return s


def txt(slide, l, t, w, h, text, pt, *, bold=False, color=C_TEXT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    font(r, pt, bold=bold, color=color)
    return box


def header_footer(slide, prs, title, subtitle, page, total):
    sw, sh = prs.slide_width, prs.slide_height
    rect(slide, 0, 0, sw, Inches(0.1), C_ACCENT)
    rect(slide, 0, Inches(0.1), sw, Inches(0.62), C_DARK)
    txt(slide, Inches(0.7), Inches(0.16), Inches(11), Inches(0.46), title, 26, bold=True, color=C_WHITE)
    if subtitle:
        txt(slide, Inches(0.7), Inches(0.64), Inches(11.2), Inches(0.38), subtitle, 15, color=RGBColor(0xB8, 0xC8, 0xD8))
    rect(slide, 0, sh - Inches(0.42), sw, Inches(0.42), RGBColor(0xF3, 0xF6, 0xF9))
    rect(slide, 0, sh - Inches(0.43), sw, Pt(1.5), C_BORDER)
    txt(slide, Inches(0.55), sh - Inches(0.34), Inches(0.6), Inches(0.24), str(page), 10, bold=True, color=C_PRIMARY)
    txt(slide, sw / 2 - Inches(0.55), sh - Inches(0.34), Inches(1.1), Inches(0.24), f"{page}/{total}", 10, color=C_MUTED, align=PP_ALIGN.CENTER)
    txt(slide, sw - Inches(3.2), sh - Inches(0.34), Inches(2.7), Inches(0.24), "FRATELANZA · Building Tomorrow Together", 9, color=C_MUTED, align=PP_ALIGN.RIGHT)


def card(slide, l, t, w, h, title, bullets, accent=C_PRIMARY):
    s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = C_CARD
    s.line.color.rgb = C_BORDER
    s.line.width = Pt(0.75)
    rect(slide, l + w - Inches(0.06), t, Inches(0.06), h, accent)
    tf = s.text_frame
    tf.margin_left = Inches(0.14)
    tf.margin_top = Inches(0.12)
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    font(r, 15, bold=True, color=C_PRIMARY)
    for b in bullets:
        p = tf.add_paragraph()
        p.space_before = Pt(4)
        p.line_spacing = 1.22
        r = p.add_run()
        r.text = f"• {b}"
        font(r, 12, color=C_TEXT)


def cards_3(slide, y, items):
    w, gap = Inches(3.85), Inches(0.35)
    x0 = Inches(0.65)
    for i, (title, bullets, accent) in enumerate(items):
        card(slide, x0 + i * (w + gap), y, w, Inches(4.75), title, bullets, accent)


def chips(slide, y, labels):
    n = len(labels)
    gap = Inches(0.18)
    total_gap = gap * (n - 1)
    w = (Inches(12.0) - total_gap) / n
    x0 = Inches(0.65)
    colors = [C_PRIMARY, C_ACCENT, C_TEAL, C_GOLD, C_PRIMARY, C_ACCENT]
    for i, label in enumerate(labels):
        s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x0 + i * (w + gap), y, w, Inches(0.58))
        s.fill.solid()
        s.fill.fore_color.rgb = colors[i % len(colors)]
        s.line.fill.background()
        tf = s.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        font(r, 12, bold=True, color=C_WHITE)


def slide_title(prs, title, subtitle, page, total, builder):
    s = blank(prs)
    header_footer(s, prs, title, subtitle, page, total)
    builder(s)
    return s


def build(prs):
    total = 18

    # 1 cover
    s = blank(prs)
    sw, sh = prs.slide_width, prs.slide_height
    rect(s, 0, 0, sw, sh, C_DARK)
    rect(s, 0, 0, sw, Inches(0.07), C_ACCENT)
    txt(s, Inches(0.8), Inches(1.8), Inches(11), Inches(0.8), "FRATELANZA", 42, bold=True, color=C_WHITE)
    txt(s, Inches(0.8), Inches(2.55), Inches(11), Inches(0.55), "Pharma Field Force Management Platform", 26, color=C_ACCENT)
    txt(s, Inches(0.8), Inches(3.15), Inches(11), Inches(0.5), "عرض فني ونطاق عمل مقترح لشركة أدوية", 18, color=RGBColor(0xCB, 0xDA, 0xEA))
    card(s, Inches(0.8), Inches(4.0), Inches(7.3), Inches(2.1), "الهدف", [
        "إدارة المندوبين والزيارات والعملاء بشكل لحظي.",
        "ربط الزيارات بالـ Sell-out والمبيعات والتحصيل.",
        "منصة قرار متكاملة — Mobile + Web + ERP.",
    ], C_ACCENT)

    # 2 agenda
    s = blank(prs)
    header_footer(s, prs, "محتويات العرض", "Technical Proposal Overview — 15 محوراً", 2, total)
    agenda = [
        ("01", "الرؤية والأهداف", "من متابعة إلى منظومة قرار"),
        ("02", "مكونات المنظومة", "Mobile · Web · Integration"),
        ("03", "إدارة Field Force", "صلاحيات · مناطق · أهداف"),
        ("04", "إدارة العملاء", "Customer 360°"),
        ("05", "Visit Management", "دورة الزيارة الكاملة"),
        ("06", "Dynamic Visit Forms", "نماذج حسب نوع العميل"),
        ("07", "Sales & Sell-out", "تحليلات قابلة للتنفيذ"),
        ("08", "Collection Tracking", "ربط التحصيل بالأداء"),
        ("09", "Control Tower", "من KPI إلى Action"),
        ("10", "Maps & Offline", "جاهز للعمل الميداني"),
        ("11", "ERP Integration", "ERP مصدر · التطبيق تشغيل"),
        ("12", "التقارير و KPIs", "مطابقة تقارير الشركة"),
        ("13", "MVP & Phase 2", "قيمة سريعة · توسع"),
        ("14", "خطة التنفيذ", "Deliverables · Acceptance"),
        ("15", "اعتماد Scope", "قبل التسعير النهائي"),
    ]
    cols = [(Inches(0.65), 0), (Inches(6.95), 8)]
    for i, (num, t, sub) in enumerate(agenda):
        col, offset = cols[0 if i < 8 else 1]
        row = i if i < 8 else i - 8
        top = Inches(1.22) + row * Inches(0.54)
        if row % 2 == 0:
            rect(s, col - Inches(0.08), top - Inches(0.04), Inches(6.0), Inches(0.48), C_ROW)
        c = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, col, top + Inches(0.06), Inches(0.4), Inches(0.4))
        c.fill.solid(); c.fill.fore_color.rgb = ACCENTS[i % 4]; c.line.fill.background()
        tf = c.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = num; font(r, 10, bold=True, color=C_WHITE)
        tx = col + Inches(0.52)
        txt(s, tx, top + Inches(0.04), Inches(5.3), Inches(0.24), t, 14, bold=True)
        txt(s, tx, top + Inches(0.28), Inches(5.3), Inches(0.22), sub, 11, color=C_MUTED)

    # content slides 3-17
    slide_title(prs, "1. الرؤية والأهداف", "من تطبيق متابعة إلى منظومة تشغيل وقرار.", 3, total, lambda s: cards_3(s, Inches(1.25), [
        ("Field Force", ["Users & Roles", "Territories", "Targets", "Route Plans", "Performance"], C_PRIMARY),
        ("Market Execution", ["Customers", "Visits", "GPS", "Dynamic Forms", "Feedback / Evidence"], C_ACCENT),
        ("Business Intelligence", ["Sales / Sell-out", "Collection / Aging", "Product Movement", "KPIs", "Action Required"], C_TEAL),
    ]))

    slide_title(prs, "2. مكونات المنظومة", "Mobile للمندوب + Web للإدارة + Integration للبيانات.", 4, total, lambda s: (
        chips(s, Inches(1.25), ["Mobile App", "Backend / API", "ERP / Data", "Analytics", "Management"]),
        cards_3(s, Inches(2.05), [
            ("Mobile App", ["Today's Visits", "Check-in / Check-out", "Visit Forms", "Photos / Tasks", "Offline + Sync"], C_ACCENT),
            ("Web Admin", ["Users", "Customers", "Products", "Targets", "Reports / Dashboard"], C_PRIMARY),
            ("Integration", ["Sales / Sell-out", "Collection", "ERP Sync", "Validation", "Export"], C_TEAL),
        ])
    ))

    slide_title(prs, "3. إدارة Field Force", "صلاحيات وهيكل جغرافي وأهداف.", 5, total, lambda s: cards_3(s, Inches(1.25), [
        ("Users & Permissions", ["Super Admin", "Management", "Regional / Area Manager", "Representative", "صلاحيات حسب الدور"], C_PRIMARY),
        ("Territories", ["Region", "Area / Territory", "Manager Assignment", "Rep Assignment", "Customer Assignment"], C_ACCENT),
        ("Targets", ["Sales Target", "Visit Target", "Collection Target", "SKU Target", "Achievement & Ranking"], C_TEAL),
    ]))

    slide_title(prs, "4. إدارة العملاء", "Customer 360° للصيدليات والأطباء والمستشفيات.", 6, total, lambda s: cards_3(s, Inches(1.25), [
        ("Customer Master", ["Pharmacy", "Doctor", "Hospital / Clinic", "Location & Territory", "Potential / Class"], C_PRIMARY),
        ("Pharma Data", ["Specialty", "Product / Brand Focus", "SKU Availability", "Competitor Products", "Competitor Price"], C_ACCENT),
        ("Customer History", ["Visits", "Sales / Sell-out", "Collection", "Feedback / Complaints", "Next Action"], C_TEAL),
    ]))

    slide_title(prs, "5. Visit Management", "الزيارة هي نقطة جمع البيانات الأساسية.", 7, total, lambda s: (
        chips(s, Inches(1.25), ["Plan", "Navigate", "Check-in", "Visit Form", "Check-out", "Sync"]),
        cards_3(s, Inches(2.05), [
            ("Planning", ["Daily / Weekly Route", "Planned vs Actual", "Unplanned Visit", "Reschedule", "Priority"], C_PRIMARY),
            ("Verification", ["GPS", "Timestamp", "Duration", "Missed Reason", "Anti-Fraud Rules"], C_ACCENT),
            ("Visit Output", ["Stock / OOS", "Feedback", "Competitors", "Complaint", "Next Action"], C_TEAL),
        ])
    ))

    slide_title(prs, "6. Dynamic Visit Forms", "Forms قابلة للتعديل حسب نوع العميل والزيارة.", 8, total, lambda s: cards_3(s, Inches(1.25), [
        ("Sales & Availability", ["Product Availability", "Stock / OOS", "Order Opportunity", "Product Focus", "Notes"], C_PRIMARY),
        ("Market Intelligence", ["Competitor", "Price", "Offer", "Display / Visibility", "Market Feedback"], C_GOLD),
        ("Evidence", ["Live Camera", "GPS + Timestamp", "Visit ID", "Attachments", "Offline + Sync"], C_TEAL),
    ]))

    slide_title(prs, "7. Sales & Sell-out Intelligence", "التقارير تتحول إلى تحليل قابل للتنفيذ.", 9, total, lambda s: cards_3(s, Inches(1.25), [
        ("Sales Views", ["By SKU", "By Brand", "By Representative", "By Territory", "By Customer"], C_PRIMARY),
        ("Product Movement", ["Top / Bottom", "Fast / Slow Moving", "Growth / Decline", "No Movement", "Out of Stock"], C_ACCENT),
        ("Comparisons", ["Current vs Previous", "Rep vs Rep", "Territory vs Territory", "Target vs Achievement", "Drill-down to SKU"], C_TEAL),
    ]))

    slide_title(prs, "8. Collection & Financial Tracking", "ربط التحصيل بأداء المندوب والعملاء.", 10, total, lambda s: cards_3(s, Inches(1.25), [
        ("Collection", ["Target", "Actual", "By Representative", "By Customer", "By Territory"], C_TEAL),
        ("Receivables", ["Outstanding", "Overdue", "Aging", "Follow-up Status", "Collection Performance"], C_GOLD),
        ("Integration", ["ERP Collection Sync", "Payment Method", "Receipt / Document", "Validation", "Excel / CSV Export"], C_PRIMARY),
    ]))

    slide_title(prs, "9. Management Control Tower", "من KPI إلى Action.", 11, total, lambda s: (
        card(s, Inches(0.65), Inches(1.25), Inches(5.9), Inches(2.35), "Live KPIs", ["Sales / Sell-out", "Visits: Planned vs Actual", "Collection / Overdue", "Achievement %"], C_PRIMARY),
        card(s, Inches(6.75), Inches(1.25), Inches(5.9), Inches(2.35), "Performance", ["Representative Ranking", "Territory Performance", "Top / Low Products", "Customer Performance"], C_ACCENT),
        card(s, Inches(0.65), Inches(3.85), Inches(5.9), Inches(2.35), "Drill-down", ["Company → Region → Territory", "→ Rep → Customer → SKU", "Filters by Date / Area / Rep"], C_TEAL),
        card(s, Inches(6.75), Inches(3.85), Inches(5.9), Inches(2.35), "Action Required", ["Declining Sales", "Missed / Low Coverage", "Slow Moving / OOS", "Overdue Collection"], C_GOLD),
    ))

    slide_title(prs, "10. Maps + Offline + Integrity", "مصمم للعمل الميداني حتى مع ضعف الإنترنت.", 12, total, lambda s: cards_3(s, Inches(1.25), [
        ("Maps & GPS", ["Customer Map", "Visit Map", "Territory Map", "Check-in / Check-out", "Planned vs Actual"], C_ACCENT),
        ("Offline-first", ["Work without Internet", "Secure Local Storage", "Background Sync", "Retry on Failure", "Sync Status"], C_TEAL),
        ("Integrity", ["GPS", "Timestamp", "User / Visit ID", "Live Camera", "Audit Trail"], C_GOLD),
    ]))

    slide_title(prs, "11. ERP & Data Integration", "ERP مصدر البيانات؛ التطبيق يضيف التشغيل والتحكم.", 13, total, lambda s: (
        chips(s, Inches(1.25), ["ERP", "API / Sync", "Validation", "Field Force", "Analytics"]),
        cards_3(s, Inches(2.05), [
            ("Data In", ["Customers", "Products / SKUs", "Sales / Sell-out", "Collection"], C_PRIMARY),
            ("Integration Rules", ["Scheduled / API Sync", "Validation", "Duplicate Prevention", "Error Log"], C_ACCENT),
            ("Data Out", ["Visits", "Feedback", "Market Intelligence", "Evidence", "Reports"], C_TEAL),
        ])
    ))

    slide_title(prs, "12. التقارير و KPIs", "التقارير الحالية للشركة هي المرجع النهائي.", 14, total, lambda s: cards_3(s, Inches(1.25), [
        ("Sell-out / Products", ["Sell-out", "Product Movement", "Product Sales", "Product Comparison", "Fast / Slow Moving"], C_PRIMARY),
        ("Field Force", ["Visits", "Coverage", "Rep Performance", "Territory Performance", "Feedback / Complaints"], C_ACCENT),
        ("Collection / Finance", ["Collection", "Outstanding", "Aging / Overdue", "Collection by Rep", "Target vs Achievement"], C_TEAL),
    ]))

    slide_title(prs, "13. MVP & Phase 2", "نبدأ بالقيمة الأساسية مع Architecture قابلة للتوسع.", 15, total, lambda s: (
        card(s, Inches(0.65), Inches(1.25), Inches(5.9), Inches(5.0), "MVP — Release 1", [
            "Users / Roles / Territories", "Customers / Doctors / Pharmacies", "Products / SKUs",
            "Route Plan + GPS Visit", "Dynamic Forms", "Sales / Sell-out + Collection",
            "Dashboard + Core Reports", "Offline + Sync",
        ], C_PRIMARY),
        card(s, Inches(6.75), Inches(1.25), Inches(5.9), Inches(5.0), "Phase 2", [
            "Live Route Tracking", "Advanced Incentives", "Advanced Alerts", "SKU / Product Targets",
            "Advanced Segmentation", "Smart Insights", "Advanced Audit / Fraud Rules", "Additional ERP Workflows",
        ], C_GOLD),
    ))

    slide_title(prs, "14. خطة التنفيذ", "كل مرحلة لها مخرجات و Acceptance Criteria.", 16, total, lambda s: [
        card(s, Inches(0.65 if i < 3 else 6.75), Inches(1.25 + (i % 3) * 1.55), Inches(5.9), Inches(1.25), t, [d], ACCENTS[i % 4])
        for i, (t, d) in enumerate([
            ("01 Discovery", "تحليل التقارير + Workshops + Final Scope"),
            ("02 UX/UI", "User Flows + Mobile + Web Prototype"),
            ("03 Architecture", "DB + APIs + Integration + Security"),
            ("04 Development", "Mobile + Web + Backend + Integration"),
            ("05 UAT", "Beta + Test Cases + Fixes"),
            ("06 Go-Live", "Deployment + Training + Handover + Warranty"),
        ])
    ] or None)

    slide_title(prs, "15. شروط اعتماد Scope", "متطلبات قبل اعتماد السعر والجدول الزمني.", 17, total, lambda s: card(s, Inches(0.65), Inches(1.25), Inches(11.95), Inches(5.0), "قبل اعتماد السعر والجدول الزمني", [
        "إرسال ومراجعة جميع التقارير الحالية للشركة.",
        "Mapping: Report → Data Source → Screen → KPI → Output.",
        "اعتماد MVP و Phase 2 كتابةً.",
        "تحديد ERP Integration ومصادر البيانات وطريقة Sync.",
        "اعتماد User Roles / Permissions و UI/UX Prototype.",
        "Acceptance Criteria لكل Milestone.",
        "Warranty / Support / Maintenance بوضوح.",
        "التسعير النهائي يُعتمد بعد إغلاق Scope.",
    ], C_ACCENT))

    # 18 closing
    s = blank(prs)
    sw, sh = prs.slide_width, prs.slide_height
    rect(s, 0, 0, sw, sh, C_DARK)
    rect(s, 0, 0, sw, Inches(0.07), C_ACCENT)
    txt(s, Inches(0.8), Inches(2.2), sw - Inches(1.6), Inches(0.65), "الخلاصة", 34, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(s, Inches(1.0), Inches(2.85), sw - Inches(2.0), Inches(0.45), "من متابعة المندوب إلى منظومة قرار متكاملة لقطاع الأدوية", 16, color=C_ACCENT, align=PP_ALIGN.CENTER)
    card(s, Inches(1.4), Inches(3.55), sw - Inches(2.8), Inches(2.1), "القيمة النهائية", [
        "رؤية موحدة لأداء المندوبين والعملاء والأصناف.",
        "بيانات موثوقة من الزيارة مع GPS و Evidence.",
        "ربط التشغيل بالمبيعات والتحصيل والـ Sell-out.",
        "Architecture قابلة للتوسع مع الـ ERP.",
    ], C_TEAL)
    txt(s, Inches(0.8), sh - Inches(1.0), sw - Inches(1.6), Inches(0.35), "شكراً لثقتكم — الخطوة التالية: Discovery Workshop", 15, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)
    txt(s, Inches(0.55), sh - Inches(0.34), Inches(0.6), Inches(0.24), "18", 10, bold=True, color=C_ACCENT)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    build(prs)
    prs.save(str(OUT))
    print(f"Built: {OUT}")


if __name__ == "__main__":
    main()
