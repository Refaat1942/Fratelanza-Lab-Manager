"""
Fix typography, spacing, and alignment for Fratelanza Pharma Field Force proposal.

Usage (Windows):
  pip install python-pptx
  python scripts/fix_pharma_pptx_typography.py

Optional:
  python scripts/fix_pharma_pptx_typography.py "C:\\path\\to\\file.pptx"
"""
from __future__ import annotations

import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# ── paths ──────────────────────────────────────────────────────────────────
DEFAULT_SRC = Path.home() / "Downloads" / "Fratelanza_Pharma_Field_Force_Technical_Proposal.pptx"
LOGO_CANDIDATES = [
    Path(r"C:\Users\a.refaat\.cursor\projects\d-Refaat-My-Projects-Fratelanza-Lab-Manager\assets")
    / "c__Users_a.refaat_AppData_Roaming_Cursor_User_workspaceStorage_empty-window_images_"
    "acfc7a18-f113-426c-b62e-74ef048d1df3-1684e138-936a-40f8-8ca0-b3e4c48459a1.png",
    Path(__file__).resolve().parent / "assets" / "fratelanza-logo.png",
]

# ── brand palette ──────────────────────────────────────────────────────────
C_PRIMARY = RGBColor(0x16, 0x69, 0xB4)
C_ACCENT = RGBColor(0x1E, 0xAE, 0xD7)
C_TEAL = RGBColor(0x19, 0x91, 0x82)
C_GOLD = RGBColor(0xEB, 0x8C, 0x2D)
C_DARK = RGBColor(0x0F, 0x2A, 0x44)
C_TEXT = RGBColor(0x1E, 0x2D, 0x3D)
C_MUTED = RGBColor(0x5C, 0x6B, 0x7A)
C_CARD = RGBColor(0xF7, 0xFA, 0xFC)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_ROW_ALT = RGBColor(0xED, 0xF2, 0xF7)
C_BORDER = RGBColor(0xD5, 0xDF, 0xEA)

FONT = "Segoe UI"

# ── type scale (pt) ───────────────────────────────────────────────────────
SZ_COVER_BRAND = 42
SZ_COVER_TITLE = 26
SZ_COVER_SUB = 18
SZ_HEADER = 26
SZ_HEADER_SUB = 15
SZ_AGENDA_TITLE = 14
SZ_AGENDA_SUB = 11
SZ_CARD_TITLE = 15
SZ_CARD_BODY = 12
SZ_STEP = 12
SZ_FOOTER = 10
SZ_FOOTER_BRAND = 9

AGENDA = [
    ("01", "الرؤية والأهداف", "من متابعة ميدانية إلى منظومة تشغيل وقرار"),
    ("02", "مكونات المنظومة", "Mobile · Web · Integration"),
    ("03", "إدارة Field Force", "صلاحيات · مناطق · أهداف"),
    ("04", "إدارة العملاء", "Customer 360° للصيدليات والأطباء"),
    ("05", "Visit Management", "دورة الزيارة الكاملة"),
    ("06", "Dynamic Visit Forms", "نماذج حسب نوع العميل والزيارة"),
    ("07", "Sales & Sell-out", "تحليلات قابلة للتنفيذ"),
    ("08", "Collection Tracking", "ربط التحصيل بأداء المندوب"),
    ("09", "Control Tower", "من KPI إلى Action"),
    ("10", "Maps & Offline", "جاهز للعمل الميداني"),
    ("11", "ERP Integration", "ERP مصدر · التطبيق طبقة تشغيل"),
    ("12", "التقارير و KPIs", "مطابقة تقارير الشركة"),
    ("13", "MVP & Phase 2", "قيمة سريعة · توسع تدريجي"),
    ("14", "خطة التنفيذ", "Deliverables · Acceptance Criteria"),
    ("15", "اعتماد Scope", "متطلبات قبل التسعير النهائي"),
]

CARD_TITLE_HINTS = (
    "Users", "Field", "Mobile", "MVP", "Live", "Discovery", "Sales", "Collection",
    "Maps", "ERP", "Sell-out", "Planning", "Verification", "Customer", "Territories",
    "Targets", "Data", "Integration", "Performance", "Drill-down", "Action", "Phase",
    "Market", "Business", "Web", "Analytics", "Management", "Receivables", "Integrity",
    "Offline", "Evidence", "Comparisons", "Product", "Pharma", "Visit", "Navigate",
    "Check-in", "Check-out", "Sync", "Plan", "القيمة", "الهدف", "قبل", "UX",
)


def find_logo() -> Path | None:
    for p in LOGO_CANDIDATES:
        if p.exists():
            return p
    return None


def rgb(c: RGBColor) -> RGBColor:
    return c


def set_font(run, size: int, *, bold=False, color: RGBColor = C_TEXT, name: str = FONT):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def remove_shape(shape):
    shape._element.getparent().remove(shape._element)


def add_rect(slide, left, top, width, height, fill: RGBColor, line: RGBColor | None = None):
    s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(0.5)
    else:
        s.line.fill.background()
    return s


def add_text(slide, left, top, width, height, text, size, *, bold=False, color=C_TEXT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    r = p.add_run()
    r.text = text
    set_font(r, size, bold=bold, color=color)
    return box


def add_header(slide, prs, title: str, subtitle: str | None = None):
    sw = prs.slide_width
    add_rect(slide, 0, 0, sw, Inches(0.1), C_ACCENT)
    add_rect(slide, 0, Inches(0.1), sw, Inches(0.62), C_DARK)
    add_text(slide, Inches(0.7), Inches(0.16), Inches(10.8), Inches(0.46), title, SZ_HEADER, bold=True, color=C_WHITE)
    if subtitle:
        add_text(slide, Inches(0.7), Inches(0.64), Inches(11.2), Inches(0.38), subtitle, SZ_HEADER_SUB, color=RGBColor(0xB8, 0xC8, 0xD8))


def add_footer(slide, prs, page: int, total: int):
    sw, sh = prs.slide_width, prs.slide_height
    add_rect(slide, 0, sh - Inches(0.42), sw, Inches(0.42), RGBColor(0xF3, 0xF6, 0xF9))
    add_rect(slide, 0, sh - Inches(0.43), sw, Pt(1.5), C_BORDER)
    add_text(slide, Inches(0.55), sh - Inches(0.34), Inches(0.7), Inches(0.24), str(page), SZ_FOOTER, bold=True, color=C_PRIMARY)
    add_text(slide, sw / 2 - Inches(0.55), sh - Inches(0.34), Inches(1.1), Inches(0.24), f"{page} / {total}", SZ_FOOTER, color=C_MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, sw - Inches(3.35), sh - Inches(0.34), Inches(2.8), Inches(0.24), "FRATELANZA · Building Tomorrow Together", SZ_FOOTER_BRAND, color=C_MUTED, align=PP_ALIGN.RIGHT)
    logo = find_logo()
    if logo:
        slide.shapes.add_picture(str(logo), sw - Inches(0.95), Inches(0.12), height=Inches(0.48))


def strip_chrome(slide):
    """Remove old headers, footers, logos, page numbers."""
    for shape in list(slide.shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            remove_shape(shape)
            continue
        if not shape.has_text_frame:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                if shape.top.inches < 0.15 and shape.height.inches < 0.25:
                    remove_shape(shape)
            continue
        text = shape.text_frame.text.strip()
        top = shape.top.inches
        if top > 6.8 and (text == "FRATELANZA" or re.fullmatch(r"\d{1,2}", text) or "/" in text or "Building Tomorrow" in text):
            remove_shape(shape)
            continue
        if top < 1.05 and (text[0:2].replace(".", "").isdigit() or text.startswith("ال") or text.startswith("محتويات") or "Overview" in text):
            remove_shape(shape)
            continue


def style_card(shape):
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.margin_left = Inches(0.14)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0.08)
    lines = tf.text.split("\n")
    tf.clear()
    for i, line in enumerate(lines):
        line = line.strip()
        if not line:
            continue
        p = tf.paragraphs[0] if i == 0 and not tf.paragraphs[0].text else tf.add_paragraph()
        p.space_before = Pt(0 if i == 0 else 3)
        p.space_after = Pt(2)
        p.line_spacing = 1.2
        is_bullet = line.startswith("•") or line.startswith("-")
        is_title = (not is_bullet) and any(line.startswith(h) for h in CARD_TITLE_HINTS)
        r = p.add_run()
        r.text = line
        if is_title:
            set_font(r, SZ_CARD_TITLE, bold=True, color=C_PRIMARY)
        else:
            set_font(r, SZ_CARD_BODY, color=C_TEXT)


def rebuild_agenda(slide, prs, page: int, total: int):
    for shape in list(slide.shapes):
        remove_shape(shape)
    add_header(slide, prs, "محتويات العرض", "Technical Proposal Overview — 15 محوراً")
    add_footer(slide, prs, page, total)

    cols = [(Inches(0.65), Inches(6.0)), (Inches(6.95), Inches(6.0))]
    row_h = Inches(0.54)
    top0 = Inches(1.22)
    badge = Inches(0.4)

    for i, (num, title, sub) in enumerate(AGENDA):
        col_idx = 0 if i < 8 else 1
        row = i if i < 8 else i - 8
        left, width = cols[col_idx]
        top = top0 + row * row_h

        if row % 2 == 0:
            add_rect(slide, left - Inches(0.08), top - Inches(0.04), width + Inches(0.16), row_h - Inches(0.06), C_ROW_ALT)

        circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left, top + Inches(0.06), badge, badge)
        circle.fill.solid()
        circle.fill.fore_color.rgb = C_PRIMARY if i % 2 == 0 else C_ACCENT
        circle.line.fill.background()
        tf = circle.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_bottom = tf.margin_top = Pt(0)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = num
        set_font(r, 10, bold=True, color=C_WHITE)

        tx = left + badge + Inches(0.14)
        tw = width - badge - Inches(0.2)
        add_text(slide, tx, top + Inches(0.04), tw, Inches(0.24), title, SZ_AGENDA_TITLE, bold=True, color=C_TEXT)
        add_text(slide, tx, top + Inches(0.28), tw, Inches(0.22), sub, SZ_AGENDA_SUB, color=C_MUTED)


def rebuild_title(slide, prs):
    sw, sh = prs.slide_width, prs.slide_height
    for shape in list(slide.shapes):
        remove_shape(shape)
    add_rect(slide, 0, 0, sw, sh, C_DARK)
    add_rect(slide, 0, 0, sw, Inches(0.07), C_ACCENT)
    add_rect(slide, 0, sh - Inches(0.07), sw, Inches(0.07), C_GOLD)
    logo = find_logo()
    if logo:
        slide.shapes.add_picture(str(logo), Inches(0.8), Inches(0.55), height=Inches(1.25))
    add_text(slide, Inches(0.8), Inches(2.05), Inches(11.5), Inches(0.75), "FRATELANZA", SZ_COVER_BRAND, bold=True, color=C_WHITE)
    add_text(slide, Inches(0.8), Inches(2.75), Inches(11.5), Inches(0.55), "Pharma Field Force Management Platform", SZ_COVER_TITLE, color=C_ACCENT)
    add_text(slide, Inches(0.8), Inches(3.35), Inches(11.5), Inches(0.5), "عرض فني ونطاق عمل مقترح لشركة أدوية", SZ_COVER_SUB, color=RGBColor(0xCB, 0xDA, 0xEA))

    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.2), Inches(7.2), Inches(2.05))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x13, 0x32, 0x50)
    card.line.color.rgb = C_ACCENT
    tf = card.text_frame
    tf.margin_left = Inches(0.22)
    tf.margin_top = Inches(0.14)
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "الهدف"
    set_font(r, 17, bold=True, color=C_ACCENT)
    for b in (
        "إدارة المندوبين والزيارات والعملاء بشكل لحظي.",
        "ربط الزيارات بالـ Sell-out والمبيعات والتحصيل.",
        "منصة قرار متكاملة — Mobile + Web + ERP.",
    ):
        p = tf.add_paragraph()
        p.space_before = Pt(5)
        p.line_spacing = 1.25
        r = p.add_run()
        r.text = f"•  {b}"
        set_font(r, 13, color=C_WHITE)

    badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.45), Inches(4.2), Inches(3.9), Inches(0.52))
    badge.fill.solid()
    badge.fill.fore_color.rgb = C_GOLD
    badge.line.fill.background()
    tf = badge.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Technical Proposal  |  Confidential"
    set_font(r, 12, bold=True, color=C_WHITE)


def rebuild_closing(slide, prs, page: int, total: int):
    sw, sh = prs.slide_width, prs.slide_height
    for shape in list(slide.shapes):
        remove_shape(shape)
    add_rect(slide, 0, 0, sw, sh, C_DARK)
    add_rect(slide, 0, 0, sw, Inches(0.07), C_ACCENT)
    logo = find_logo()
    if logo:
        slide.shapes.add_picture(str(logo), sw / 2 - Inches(0.7), Inches(0.95), height=Inches(1.35))
    add_text(slide, Inches(0.8), Inches(2.65), sw - Inches(1.6), Inches(0.65), "الخلاصة", 34, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1.0), Inches(3.25), sw - Inches(2.0), Inches(0.45), "من متابعة المندوب إلى منظومة قرار متكاملة لقطاع الأدوية", 16, color=C_ACCENT, align=PP_ALIGN.CENTER)

    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.4), Inches(4.0), sw - Inches(2.8), Inches(2.05))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x13, 0x32, 0x50)
    card.line.color.rgb = C_TEAL
    tf = card.text_frame
    tf.margin_left = Inches(0.25)
    tf.margin_top = Inches(0.12)
    tf.clear()
    for j, b in enumerate((
        "رؤية موحدة لأداء المندوبين والعملاء والأصناف.",
        "بيانات موثوقة من الزيارة مع GPS و Evidence.",
        "ربط التشغيل بالمبيعات والتحصيل والـ Sell-out.",
        "Architecture قابلة للتوسع مع الـ ERP.",
    )):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.space_before = Pt(4)
        p.line_spacing = 1.25
        r = p.add_run()
        r.text = f"•  {b}"
        set_font(r, 13, color=C_WHITE)

    add_text(slide, Inches(0.8), sh - Inches(1.0), sw - Inches(1.6), Inches(0.35), "شكراً لثقتكم — الخطوة التالية: Discovery Workshop", 15, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.8), sh - Inches(0.58), sw - Inches(1.6), Inches(0.28), "FRATELANZA · info@fratelanza.com", 11, color=C_MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.55), sh - Inches(0.34), Inches(0.7), Inches(0.24), str(page), SZ_FOOTER, bold=True, color=C_ACCENT)
    add_text(slide, sw / 2 - Inches(0.55), sh - Inches(0.34), Inches(1.1), Inches(0.24), f"{page} / {total}", SZ_FOOTER, color=C_MUTED, align=PP_ALIGN.CENTER)


def detect_title_subtitle(slide) -> tuple[str | None, str | None]:
    title = subtitle = None
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.text.strip()
        top = shape.top.inches
        if top > 1.1 or not t:
            continue
        if title is None and (re.match(r"^\d+\.", t) or t.startswith("ال") or "Scope" in t):
            title = t
        elif subtitle is None and title and t != title:
            subtitle = t
    return title, subtitle


def fix_content_slide(slide, prs, page: int, total: int):
    title, subtitle = detect_title_subtitle(slide)
    strip_chrome(slide)
    if title:
        add_header(slide, prs, title, subtitle)
    add_footer(slide, prs, page, total)
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and shape.has_text_frame:
            if shape.top.inches > 1.15:
                style_card(shape)
                shape.line.color.rgb = C_BORDER
                shape.line.width = Pt(0.75)


def is_agenda(slide) -> bool:
    for shape in slide.shapes:
        if shape.has_text_frame and "محتويات العرض" in shape.text_frame.text:
            return True
    return False


def main():
    src = Path(sys.argv[1]) if len(sys.argv) > 1 else DEFAULT_SRC
    if not src.exists():
        print(f"File not found: {src}")
        print("Pass path as argument: python fix_pharma_pptx_typography.py \"C:\\\\path\\\\to\\\\file.pptx\"")
        sys.exit(1)

    out = src.with_name(src.stem + "_FIXED.pptx")
    prs = Presentation(str(src))
    total = len(prs.slides)

    rebuild_title(prs.slides[0], prs)

    for idx, slide in enumerate(prs.slides):
        if idx == 0:
            continue
        page = idx + 1
        if idx == 1 or is_agenda(slide):
            rebuild_agenda(slide, prs, page, total)
        elif idx == total - 1:
            rebuild_closing(slide, prs, page, total)
        else:
            fix_content_slide(slide, prs, page, total)

    prs.save(str(out))
    print(f"Saved: {out}")

    # overwrite source for convenience
    import shutil
    bak = src.with_suffix(".pptx.bak2")
    if not bak.exists():
        shutil.copy2(src, bak)
    shutil.copy2(out, src)
    print(f"Updated: {src}")


if __name__ == "__main__":
    main()
